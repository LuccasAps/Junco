"""
Gera LogBI — Lead Time Analytics.pptx e .pdf
Cada slide: header de marca + imagem do dashboard + painel lateral com insights.
"""
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib import colors
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

IMG_DIR  = Path("imagens")
OUT_PPTX = Path("LogBI_LeadTime_Analytics.pptx")
OUT_PDF  = Path("LogBI_LeadTime_Analytics.pdf")

# ── Paleta ──────────────────────────────────────────────────────────────────
C_DARK  = RGBColor(0x1E, 0x3A, 0x5F)
C_BLUE  = RGBColor(0x25, 0x63, 0xEB)
C_LIGHT = RGBColor(0xF0, 0xF4, 0xFF)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY  = RGBColor(0x64, 0x74, 0x8B)
C_RED   = RGBColor(0xDC, 0x26, 0x26)
C_GREEN = RGBColor(0x16, 0xA3, 0x4A)

# ── Dimensões do slide 16:9 ──────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)

# ── Conteúdo dos slides ──────────────────────────────────────────────────────
# (arquivo, título, subtítulo, lista de bullets)
SLIDES = [
    (
        "indicadores principais.png",
        "Indicadores Principais",
        "Visão consolidada do desempenho logístico no período",
        [
            "Lead Time Total médio de 9,9 dias (Pedido→Entrega)",
            "NF→Expedição: 4,2 dias em média",
            "Expedição→Entrega: 6,1 dias em média",
            "Apenas 3,5% das 21.185 entregas foram no prazo",
            "20.446 ocorrências de atraso registradas",
        ],
    ),
    (
        "matriz de intervalos.png",
        "Matrix de Intervalos por Etapa",
        "Decomposição do Lead Time em cada fase do processo",
        [
            "Pedido→NF: 0,0 dias — emissão instantânea",
            "NF→Expedição: 4,2 dias (22% acima da média)",
            "Expedição→Entrega: 6,1 dias (36% acima da média)",
            "NF→Entrega: 9,9 dias (40% acima da média)",
            "Etapa de transporte é o principal gargalo",
        ],
    ),
    (
        "tendencia de lead time.png",
        "Tendência de Lead Time — Mensal",
        "Evolução do Lead Time de janeiro a fevereiro de 2026",
        [
            "Lead Time Total caiu de ~10,2 para ~9,3 dias",
            "Expedição→Entrega estável em ~6 dias",
            "Pedido→NF manteve-se em 0 dias no período",
            "Tendência de queda indica melhoria em andamento",
        ],
    ),
    (
        "Pontualidade de entregas.png",
        "Pontualidade de Entregas",
        "Proporção de entregas realizadas no prazo vs. em atraso",
        [
            "96,5% das entregas realizadas em atraso",
            "Apenas 3,5% concluídas dentro do prazo",
            "Situação crítica exige ação imediata",
            "Benchmark do setor: mínimo 85% no prazo",
            "Oportunidade de melhoria significativa",
        ],
    ),
    (
        "Análise por dimensão.png",
        "Análise por Dimensão",
        "Lead Time segmentado por UF, transportadora e região",
        [
            "Nordeste: maior Lead Time regional (18 dias)",
            "Norte: 16,8 dias — segunda região mais crítica",
            "Sudeste: melhor desempenho (7,8 dias)",
            "RR e AP: UFs com maior Lead Time (30,5 e 29 dias)",
            "ACM Soluções Logística: pior transportadora (18,2 d)",
        ],
    ),
    (
        "Ranking de Transportadoras.png",
        "Ranking de Transportadoras",
        "Melhores e piores parceiros logísticos por Lead Time",
        [
            "Top 3 melhores: Lead Time abaixo de 0,6 dias",
            "ACM Soluções: 18,2 dias e apenas 0,6% no prazo",
            "EVIDENCIA Logística: 16,5 dias, 640 entregas",
            "Via Connect: 14,4 dias, 46,3% de pontualidade",
            "JCL Transportes: marcada como atenção (12,8 d)",
        ],
    ),
    (
        "Tabela Analitica.png",
        "Tabela Analítica Detalhada",
        "Exploração granular nota a nota com filtros dinâmicos",
        [
            "Visualização individual por nota fiscal",
            "Campos: emissão, expedição, entrega, cliente, UF",
            "Filtros dinâmicos por qualquer dimensão",
            "Permite auditoria de casos específicos",
            "21.185 registros disponíveis para análise",
        ],
    ),
    (
        "Filtros.png",
        "Painel de Filtros",
        "Controles interativos para segmentação dos dados",
        [
            "Filtro por UF de destino",
            "Filtro por transportadora",
            "Filtro por cliente e operação fiscal",
            "Filtro por região geográfica",
            "Seleção de período por emissão da NF",
        ],
    ),
]

# ────────────────────────────────────────────────────────────────────────────
# Helpers PPTX
# ────────────────────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, rgb, alpha=None):
    from pptx.util import Pt as _Pt
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = rgb
    s.line.fill.background()
    return s


def txb(slide, text, l, t, w, h, size, bold=False,
        rgb=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = rgb
    return box


def bullet_box(slide, items, l, t, w, h, size=11, color=C_DARK):
    from pptx.util import Pt as _Pt
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = f"• {item}"
        r.font.size  = Pt(size)
        r.font.color.rgb = color
    return box


def fit_image(slide, img_path, l, t, max_w, max_h):
    img = Image.open(img_path)
    iw, ih = img.size
    scale = min(max_w / iw, max_h / ih)
    w = Emu(int(iw * scale))
    h = Emu(int(ih * scale))
    cx = l + (max_w - w) // 2
    cy = t + (max_h - h) // 2
    slide.shapes.add_picture(str(img_path), cx, cy, w, h)


# ────────────────────────────────────────────────────────────────────────────
# Slides PPTX
# ────────────────────────────────────────────────────────────────────────────

def slide_cover(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    # fundo azul escuro total
    rect(sl, 0, 0, SW, SH, C_DARK)

    # faixa lateral direita azul vivo
    rect(sl, Inches(8.8), 0, Inches(4.53), SH, C_BLUE)

    # linha decorativa
    rect(sl, Inches(8.65), 0, Inches(0.06), SH, C_WHITE)

    # logotipo / título
    txb(sl, "LogBI", Inches(0.7), Inches(1.6), Inches(7.5), Inches(1.8),
        72, bold=True, rgb=C_WHITE)
    txb(sl, "Lead Time Analytics", Inches(0.7), Inches(3.3), Inches(7.5), Inches(0.8),
        28, rgb=RGBColor(0xBF, 0xD7, 0xFF))
    txb(sl, "Dashboard de desempenho logístico completo\n"
            "Lead Time · Pontualidade · Transportadoras · Regiões",
        Inches(0.7), Inches(4.2), Inches(7.5), Inches(1.0),
        14, rgb=RGBColor(0xA0, 0xB8, 0xE8))
    txb(sl, "Março 2026", Inches(0.7), Inches(6.4), Inches(4), Inches(0.5),
        12, rgb=RGBColor(0x80, 0xA0, 0xCC))

    # painel direito — destaques
    txb(sl, "Destaques", Inches(9.1), Inches(0.8), Inches(3.8), Inches(0.5),
        16, bold=True, rgb=C_WHITE)
    highlights = [
        "21.185 entregas analisadas",
        "Lead Time médio: 9,9 dias",
        "96,5% das entregas em atraso",
        "8 seções de análise interativa",
    ]
    bullet_box(sl, highlights, Inches(9.1), Inches(1.5), Inches(3.9), Inches(3),
               size=13, color=C_WHITE)


def slide_agenda(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, SW, Inches(1.1), C_DARK)
    rect(sl, 0, Inches(1.1), SW, SH - Inches(1.1), C_LIGHT)
    rect(sl, 0, SH - Inches(0.35), SW, Inches(0.35), C_BLUE)

    txb(sl, "Agenda", Inches(0.4), Inches(0.2), Inches(10), Inches(0.7),
        24, bold=True, rgb=C_WHITE)
    txb(sl, "Conteúdo desta apresentação", Inches(0.4), Inches(6.85),
        Inches(10), Inches(0.3), 9, rgb=C_WHITE)

    agenda_items = [
        ("01", "Indicadores Principais", "Resumo executivo do desempenho logístico"),
        ("02", "Matrix de Intervalos",   "Decomposição do Lead Time por etapa"),
        ("03", "Tendência Mensal",       "Evolução ao longo do tempo"),
        ("04", "Pontualidade",           "Análise de entregas no prazo vs. atraso"),
        ("05", "Análise por Dimensão",   "Segmentação por UF, transportadora e região"),
        ("06", "Ranking de Transportadoras", "Melhores e piores parceiros"),
        ("07", "Tabela Analítica",       "Exploração detalhada nota a nota"),
        ("08", "Painel de Filtros",      "Controles de segmentação do dashboard"),
    ]

    cols = [Inches(0.5), Inches(6.9)]
    for i, (num, title, desc) in enumerate(agenda_items):
        col = cols[i % 2]
        row = i // 2
        top = Inches(1.4) + row * Inches(1.2)

        rect(sl, col, top, Inches(0.55), Inches(0.55), C_BLUE)
        txb(sl, num, col + Inches(0.05), top + Inches(0.05),
            Inches(0.45), Inches(0.45), 14, bold=True, rgb=C_WHITE, align=PP_ALIGN.CENTER)
        txb(sl, title, col + Inches(0.65), top, Inches(5.8), Inches(0.35),
            13, bold=True, rgb=C_DARK)
        txb(sl, desc, col + Inches(0.65), top + Inches(0.35), Inches(5.8), Inches(0.35),
            10, rgb=C_GRAY)


def slide_content(prs, img_path, title, subtitle, bullets):
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Header ────────────────────────────────────────────────────────────
    rect(sl, 0, 0, SW, Inches(1.0), C_DARK)
    rect(sl, 0, Inches(1.0), SW, Inches(0.04), C_BLUE)
    txb(sl, title, Inches(0.4), Inches(0.1), Inches(9), Inches(0.55),
        22, bold=True, rgb=C_WHITE)
    txb(sl, subtitle, Inches(0.4), Inches(0.65), Inches(9.5), Inches(0.32),
        11, rgb=RGBColor(0xBF, 0xD7, 0xFF))

    # ── Rodapé ─────────────────────────────────────────────────────────────
    rect(sl, 0, SH - Inches(0.32), SW, Inches(0.32), C_BLUE)
    txb(sl, "LogBI — Lead Time Analytics  |  Confidencial",
        Inches(0.3), SH - Inches(0.3), Inches(8), Inches(0.28), 8, rgb=C_WHITE)
    txb(sl, "Março 2026",
        Inches(11.5), SH - Inches(0.3), Inches(1.6), Inches(0.28),
        8, rgb=C_WHITE, align=PP_ALIGN.RIGHT)

    # ── Painel lateral de insights ─────────────────────────────────────────
    panel_l = Inches(9.85)
    panel_t = Inches(1.1)
    panel_w = Inches(3.28)
    panel_h = SH - Inches(1.42)

    rect(sl, panel_l, panel_t, panel_w, panel_h, C_LIGHT)
    rect(sl, panel_l, panel_t, Inches(0.06), panel_h, C_BLUE)

    txb(sl, "Principais Insights",
        panel_l + Inches(0.15), panel_t + Inches(0.15),
        panel_w - Inches(0.2), Inches(0.4),
        12, bold=True, rgb=C_DARK)

    bullet_box(sl, bullets,
               panel_l + Inches(0.15), panel_t + Inches(0.65),
               panel_w - Inches(0.25), panel_h - Inches(0.8),
               size=10, color=C_DARK)

    # ── Imagem do dashboard ────────────────────────────────────────────────
    img_area_l = Inches(0.2)
    img_area_t = Inches(1.1)
    img_area_w = Inches(9.5)
    img_area_h = SH - Inches(1.45)

    fit_image(sl, img_path, img_area_l, img_area_t, img_area_w, img_area_h)


def build_pptx():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    slide_cover(prs)
    slide_agenda(prs)
    for fname, title, subtitle, bullets in SLIDES:
        slide_content(prs, IMG_DIR / fname, title, subtitle, bullets)

    prs.save(OUT_PPTX)
    print(f"PPTX salvo: {OUT_PPTX}")


# ────────────────────────────────────────────────────────────────────────────
# PDF — uma página por slide, layout fiel ao PPTX
# ────────────────────────────────────────────────────────────────────────────

PW, PH = landscape(A4)   # 841.89 x 595.28 pt

def hex2rgb_f(h):
    r = int(h[1:3], 16) / 255
    g = int(h[3:5], 16) / 255
    b = int(h[5:7], 16) / 255
    return r, g, b

DARK_F  = hex2rgb_f("#1E3A5F")
BLUE_F  = hex2rgb_f("#2563EB")
LIGHT_F = hex2rgb_f("#F0F4FF")
WHITE_F = (1, 1, 1)
GRAY_F  = hex2rgb_f("#64748B")

def pdf_rect(c, x, y, w, h, fill, stroke=False):
    c.setFillColorRGB(*fill)
    if stroke:
        c.setStrokeColorRGB(*fill)
        c.rect(x, y, w, h, fill=1, stroke=1)
    else:
        c.rect(x, y, w, h, fill=1, stroke=0)

def pdf_text(c, text, x, y, size, bold=False, color=WHITE_F, align="left"):
    c.setFillColorRGB(*color)
    c.setFont("Helvetica-Bold" if bold else "Helvetica", size)
    if align == "center":
        c.drawCentredString(x, y, text)
    elif align == "right":
        c.drawRightString(x, y, text)
    else:
        c.drawString(x, y, text)

def pdf_fit_image(c, img_path, x, y, max_w, max_h):
    img = Image.open(img_path)
    iw, ih = img.size
    scale = min(max_w / iw, max_h / ih)
    w = iw * scale
    h = ih * scale
    cx = x + (max_w - w) / 2
    cy = y + (max_h - h) / 2
    c.drawImage(str(img_path), cx, cy, w, h, preserveAspectRatio=True)

def pdf_cover(c):
    # fundo
    pdf_rect(c, 0, 0, PW, PH, DARK_F)
    # faixa direita
    split = PW * 0.66
    pdf_rect(c, split, 0, PW - split, PH, BLUE_F)
    # linha separadora
    pdf_rect(c, split - 2, 0, 4, PH, WHITE_F)

    # título
    pdf_text(c, "LogBI", 40, PH - 120, 60, bold=True)
    pdf_text(c, "Lead Time Analytics", 40, PH - 160, 22, color=hex2rgb_f("#BFCFFF"))
    pdf_text(c, "Dashboard de desempenho logístico completo", 40, PH - 190, 12, color=hex2rgb_f("#A0B8E8"))
    pdf_text(c, "Lead Time  ·  Pontualidade  ·  Transportadoras  ·  Regiões", 40, PH - 207, 12, color=hex2rgb_f("#A0B8E8"))
    pdf_text(c, "Março 2026", 40, 40, 11, color=hex2rgb_f("#80A0CC"))

    # painel direito
    rx = split + 20
    pdf_text(c, "Destaques", rx, PH - 90, 16, bold=True)
    highlights = [
        "21.185 entregas analisadas",
        "Lead Time médio: 9,9 dias",
        "96,5% das entregas em atraso",
        "8 seções de análise interativa",
    ]
    for i, h in enumerate(highlights):
        pdf_text(c, f"•  {h}", rx, PH - 130 - i * 24, 12)

def pdf_content_page(c, img_path, title, subtitle, bullets):
    HEADER_H = 60
    FOOTER_H = 20
    PANEL_W  = PW * 0.24
    IMG_PAD  = 6

    # fundo branco
    pdf_rect(c, 0, 0, PW, PH, WHITE_F)

    # header
    pdf_rect(c, 0, PH - HEADER_H, PW, HEADER_H, DARK_F)
    pdf_rect(c, 0, PH - HEADER_H - 3, PW, 3, BLUE_F)
    pdf_text(c, title, 14, PH - 36, 18, bold=True)
    pdf_text(c, subtitle, 14, PH - 55, 9, color=hex2rgb_f("#BFCFFF"))

    # rodapé
    pdf_rect(c, 0, 0, PW, FOOTER_H, BLUE_F)
    pdf_text(c, "LogBI — Lead Time Analytics  |  Confidencial", 10, 5, 7)
    pdf_text(c, "Março 2026", PW - 10, 5, 7, align="right")

    # painel lateral insights
    panel_x = PW - PANEL_W - 6
    panel_y = FOOTER_H + 4
    panel_h = PH - HEADER_H - FOOTER_H - 10
    pdf_rect(c, panel_x, panel_y, PANEL_W, panel_h, LIGHT_F)
    pdf_rect(c, panel_x, panel_y, 4, panel_h, BLUE_F)

    c.setFillColorRGB(*DARK_F)
    c.setFont("Helvetica-Bold", 9)
    c.drawString(panel_x + 10, panel_y + panel_h - 20, "Principais Insights")

    c.setFont("Helvetica", 8)
    y_b = panel_y + panel_h - 38
    for bullet in bullets:
        # quebra de linha simples para textos longos
        words = bullet.split()
        line = ""
        for word in words:
            test = f"• {line} {word}".strip() if line else f"• {word}"
            if c.stringWidth(test, "Helvetica", 8) < PANEL_W - 18:
                line = (line + " " + word).strip()
            else:
                if y_b < panel_y + 8:
                    break
                c.setFillColorRGB(*DARK_F)
                c.drawString(panel_x + 10, y_b, f"• {line}")
                y_b -= 13
                line = word
        if line and y_b >= panel_y + 8:
            c.setFillColorRGB(*DARK_F)
            c.drawString(panel_x + 10, y_b, f"• {line}")
            y_b -= 18

    # imagem
    img_x = IMG_PAD
    img_y = FOOTER_H + 4
    img_w = panel_x - IMG_PAD * 2
    img_h = PH - HEADER_H - FOOTER_H - 10
    pdf_fit_image(c, img_path, img_x, img_y, img_w, img_h)


def build_pdf():
    c = rl_canvas.Canvas(str(OUT_PDF), pagesize=(PW, PH))

    pdf_cover(c)
    c.showPage()

    for fname, title, subtitle, bullets in SLIDES:
        pdf_content_page(c, IMG_DIR / fname, title, subtitle, bullets)
        c.showPage()

    c.save()
    print(f"PDF salvo:  {OUT_PDF}")


if __name__ == "__main__":
    build_pptx()
    build_pdf()